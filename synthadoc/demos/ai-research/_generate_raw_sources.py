"""Generate binary raw_sources for the ai-research demo.

Run once from this directory:
    python _generate_raw_sources.py
"""
from pathlib import Path
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from PIL import Image, ImageDraw, ImageFont

OUT = Path(__file__).parent / "raw_sources"


# ── 1. model-capabilities-comparison.xlsx ────────────────────────────────────

def make_xlsx():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Model Comparison"

    header_fill = PatternFill("solid", fgColor="1F4E79")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    alt_fill = PatternFill("solid", fgColor="D6E4F0")

    headers = [
        "Model", "Organisation", "Release", "Parameters",
        "MMLU (5-shot)", "HumanEval", "Context Window", "Open Weights",
    ]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")

    rows = [
        ("GPT-3",         "OpenAI",    "2020-06", "175B",       "~43%", "—",     "4K",   "No"),
        ("BERT-Large",    "Google",    "2018-10", "340M",       "—",    "—",     "512",  "Yes"),
        ("GPT-4",         "OpenAI",    "2023-03", "Undisclosed","86.4%","67.0%", "128K", "No"),
        ("LLaMA 2 70B",   "Meta",      "2023-07", "70B",        "68.9%","29.9%", "4K",   "Yes"),
        ("Llama 3 70B",   "Meta",      "2024-04", "70B",        "82.0%","81.7%", "8K",   "Yes"),
        ("Claude 3 Opus", "Anthropic", "2024-03", "Undisclosed","86.8%","84.9%", "200K", "No"),
        ("Gemini Ultra",  "Google",    "2023-12", "Undisclosed","90.0%","74.4%", "32K",  "No"),
        ("Mistral 7B",    "Mistral",   "2023-09", "7B",         "60.1%","30.5%", "8K",   "Yes"),
    ]
    for r, row in enumerate(rows, 2):
        fill = alt_fill if r % 2 == 0 else None
        for c, val in enumerate(row, 1):
            cell = ws.cell(row=r, column=c, value=val)
            if fill:
                cell.fill = fill
            cell.alignment = Alignment(horizontal="center")

    ws.column_dimensions["A"].width = 18
    ws.column_dimensions["B"].width = 14
    for col in "CDEFGH":
        ws.column_dimensions[col].width = 13

    # Notes sheet
    ws2 = wb.create_sheet("Notes")
    ws2["A1"] = "Benchmark Notes"
    ws2["A1"].font = Font(bold=True, size=12)
    notes = [
        ("MMLU", "Massive Multitask Language Understanding — 57 subjects, 5-shot unless noted."),
        ("HumanEval", "Python code generation — pass@1 rate on 164 problems."),
        ("Gemini Ultra MMLU", "The 90.0% figure uses CoT@32 (32 samples with majority vote). "
                              "Standard 5-shot gives 83.7% — below GPT-4's 86.4% on the same protocol."),
        ("Context Window", "As of initial release; many models have since expanded via API updates."),
    ]
    for r, (term, note) in enumerate(notes, 3):
        ws2.cell(row=r, column=1, value=term).font = Font(bold=True)
        ws2.cell(row=r, column=2, value=note)
    ws2.column_dimensions["A"].width = 22
    ws2.column_dimensions["B"].width = 80

    path = OUT / "model-capabilities-comparison.xlsx"
    wb.save(path)
    print(f"  wrote {path.name} ({path.stat().st_size:,} bytes)")


# ── 2. deep-learning-concepts.pptx ───────────────────────────────────────────

def make_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]   # blank
    title_content = prs.slide_layouts[1]

    def add_title_slide(title, subtitle=""):
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = title
        sl.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
        sl.placeholders[1].text = subtitle

    def add_content_slide(title, bullets):
        sl = prs.slides.add_slide(title_content)
        sl.shapes.title.text = title
        tf = sl.placeholders[1].text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

    add_title_slide(
        "Deep Learning Concepts",
        "A concise reference for AI/ML practitioners",
    )
    add_content_slide("The Neural Network Stack", [
        "Input layer → hidden layers → output layer",
        "Each layer applies: linear transform + non-linear activation",
        "Activation functions: ReLU, GELU, SiLU — introduce non-linearity",
        "Training: forward pass (compute loss) + backward pass (update weights via gradient descent)",
        "Optimisers: SGD, Adam, AdamW — AdamW standard for LLM training",
    ])
    add_content_slide("The Transformer in Brief", [
        "Introduced by Vaswani et al. (2017) — 'Attention Is All You Need'",
        "Replaces recurrent layers with multi-head self-attention",
        "Self-attention: every token attends to every other token — O(N²) per layer",
        "Positional encodings added to embeddings (sinusoidal or learned)",
        "Pre-norm (LayerNorm before attention) is now standard for stability",
    ])
    add_content_slide("Attention Mechanism", [
        "Query (Q), Key (K), Value (V) projections from input",
        "Attention scores: softmax(QK^T / sqrt(d_k)) · V",
        "Multi-head: run h parallel attention heads, concatenate outputs",
        "Each head learns different relationship patterns in the data",
        "Flash Attention: exact same result, O(N) memory via kernel fusion",
    ])
    add_content_slide("Training Pipeline", [
        "1. Pre-training: next-token prediction on ~1T–15T tokens",
        "2. Supervised fine-tuning (SFT): instruction-following examples",
        "3. Reward model: human rankings of output quality",
        "4. RLHF/DPO: align model to human preferences",
        "Compute budget determines model size and token count (Chinchilla laws)",
    ])
    add_content_slide("Key Scaling Findings", [
        "Loss follows a power law in N (params), D (tokens), C (compute)",
        "Kaplan et al. 2020: model size matters most per unit compute",
        "Chinchilla 2022: tokens should scale ~20x parameters",
        "GPT-3 / PaLM were significantly under-trained by Chinchilla standard",
        "LLaMA series designed to be compute-optimal or intentionally over-trained",
    ])

    path = OUT / "deep-learning-concepts.pptx"
    prs.save(path)
    print(f"  wrote {path.name} ({path.stat().st_size:,} bytes)")


# ── 3. llm-benchmarks-q1-2026.pdf ────────────────────────────────────────────

def make_pdf():
    path = OUT / "llm-benchmarks-q1-2026.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    W, H = letter

    def line(y, text, size=11, bold=False, indent=0):
        c.setFont("Helvetica-Bold" if bold else "Helvetica", size)
        c.drawString(72 + indent, y, text)

    y = H - 72
    line(y, "Gemini Ultra MMLU: Correcting a Widely Cited Error", size=16, bold=True); y -= 32
    line(y, "AI Benchmark Watch — Q1 2026", size=10); y -= 40

    for text in [
        "A persistent error in LLM benchmark reporting claims that Gemini Ultra was the first",
        "model to surpass human expert performance on MMLU, citing its 90.0% score against the",
        "89.8% human expert baseline. This claim is wrong.",
        "",
        "The 90.0% figure was produced using CoT@32 evaluation — 32 chain-of-thought samples",
        "per question with majority voting. Every other model in the standard comparison table",
        "(GPT-4 at 86.4%, Claude 3 Opus at 86.8%) was evaluated under the standard 5-shot",
        "direct-answer protocol. These are not comparable numbers.",
        "",
        "Under the equivalent 5-shot protocol, Gemini Ultra scores 83.7% — below GPT-4.",
        "No frontier model has surpassed human expert performance on MMLU under the standard",
        "5-shot evaluation. The claim that Gemini Ultra was the first to do so is false.",
        "",
        "The error originates in Google's December 2023 technical report, which listed the",
        "CoT@32 result alongside 5-shot results from competing models without clearly",
        "distinguishing the evaluation protocols. Subsequent coverage repeated the mistake.",
        "",
        "The llm-benchmarks wiki entry repeats this error in its results table and explicitly",
        "states that Gemini Ultra's 90.0% 'is the first result to surpass human expert",
        "performance (89.8%) on this benchmark, marking a significant milestone.' This claim",
        "disputes the standard benchmark comparison and should be marked as contradicted.",
    ]:
        line(y, text); y -= 16

    c.save()
    print(f"  wrote {path.name} ({path.stat().st_size:,} bytes)")


# ── 4. neural-network-architecture.png ───────────────────────────────────────

def make_png():
    W, H = 900, 600
    img = Image.new("RGB", (W, H), "#1a1a2e")
    draw = ImageDraw.Draw(img)

    title_color = "#e0e0e0"
    node_colors = {"input": "#4a90d9", "hidden": "#7b68ee", "output": "#50c878"}
    edge_color = "#555577"
    label_color = "#cccccc"

    # Title
    draw.text((W // 2, 30), "Neural Network Architecture", fill=title_color, anchor="mm")
    draw.text((W // 2, 55), "Feedforward Network — Fully Connected Layers", fill="#888888", anchor="mm")

    # Layer positions
    layers = [
        ("Input\nLayer",  "input",  [150]),
        ("Hidden\nLayer 1", "hidden", [250]),
        ("Hidden\nLayer 2", "hidden", [350]),
        ("Output\nLayer", "output", [550]),  # using x as position holder
    ]

    # x positions and node counts
    x_positions = [120, 270, 420, 600, 750]
    node_counts  = [4, 6, 6, 5, 3]
    layer_names  = ["Input", "Hidden 1", "Hidden 2", "Hidden 3", "Output"]
    layer_types  = ["input", "hidden", "hidden", "hidden", "output"]

    node_radius = 18
    node_positions = []

    for li, (x, n, ltype) in enumerate(zip(x_positions, node_counts, layer_types)):
        y_start = (H - n * 70) // 2 + 60
        layer_nodes = []
        for ni in range(n):
            y = y_start + ni * 70
            layer_nodes.append((x, y))
        node_positions.append(layer_nodes)

    # Draw edges (previous layer to current)
    for li in range(1, len(node_positions)):
        for nx, ny in node_positions[li]:
            for px, py in node_positions[li - 1]:
                draw.line([(px + node_radius, py), (nx - node_radius, ny)], fill=edge_color, width=1)

    # Draw nodes
    for li, (layer_nodes, ltype) in enumerate(zip(node_positions, layer_types)):
        color = node_colors.get(ltype, "#7b68ee")
        for nx, ny in layer_nodes:
            draw.ellipse(
                [(nx - node_radius, ny - node_radius), (nx + node_radius, ny + node_radius)],
                fill=color, outline="#ffffff", width=2,
            )

    # Layer labels
    for li, (x, name) in enumerate(zip(x_positions, layer_names)):
        draw.text((x, H - 35), name, fill=label_color, anchor="mm")

    # Legend
    for i, (label, color) in enumerate([("Input", "#4a90d9"), ("Hidden", "#7b68ee"), ("Output", "#50c878")]):
        lx, ly = 30, 120 + i * 35
        draw.ellipse([(lx, ly - 8), (lx + 16, ly + 8)], fill=color)
        draw.text((lx + 24, ly), label, fill=label_color, anchor="lm")

    draw.text((W // 2, H - 10),
              "Diagram: nodes = neurons, edges = weighted connections, colour = layer type",
              fill="#666666", anchor="mm")

    path = OUT / "neural-network-architecture.png"
    img.save(path, "PNG")
    print(f"  wrote {path.name} ({path.stat().st_size:,} bytes)")


if __name__ == "__main__":
    print("Generating ai-research raw_sources binary files...")
    make_xlsx()
    make_pptx()
    make_pdf()
    make_png()
    print("Done.")
